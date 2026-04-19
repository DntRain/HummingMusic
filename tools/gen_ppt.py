"""
tools/gen_ppt.py - 生成中期答辩 PPT
用法：python tools/gen_ppt.py
输出：reports/midterm_xyk.pptx
"""

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
import numpy as np

# 注册微软雅黑用于图表中文
_MSYH = "/usr/share/fonts/WindowsFonts/msyh.ttc"
_MSYH_BD = "/usr/share/fonts/WindowsFonts/msyhbd.ttc"
for _fp in [_MSYH, _MSYH_BD]:
    if Path(_fp).exists():
        fm.fontManager.addfont(_fp)
plt.rcParams["font.family"] = "Microsoft YaHei"
plt.rcParams["axes.unicode_minus"] = False

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 颜色 ──────────────────────────────────────
C_DARK   = RGBColor(0x1F, 0x38, 0x64)   # 深蓝
C_MID    = RGBColor(0x2E, 0x75, 0xB6)   # 中蓝
C_ACCENT = RGBColor(0xED, 0x7D, 0x31)   # 橙色（高亮你的模块）
C_GREEN  = RGBColor(0x70, 0xAD, 0x47)   # 绿色（达标）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xD6, 0xE4, 0xF0)   # 浅蓝背景
C_GRAY   = RGBColor(0x59, 0x59, 0x59)
C_BG     = RGBColor(0xF5, 0xF8, 0xFC)   # 页面背景

FONT = "微软雅黑"
W, H = Inches(13.33), Inches(7.5)   # 16:9


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 完全空白
    return prs.slides.add_slide(layout)


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_para(tf, text, size=Pt(16), bold=False, color=None,
             align=PP_ALIGN.LEFT, indent=0, italic=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def slide_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def header_bar(slide, title, subtitle=None):
    """顶部深蓝色标题栏"""
    bar = add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill=C_DARK)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), Inches(10), Inches(0.85),
             size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(10.5), Inches(0.28), Inches(2.5), Inches(0.6),
                 size=Pt(14), color=C_LIGHT, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════
def slide_cover(prs):
    slide = blank_slide(prs)
    slide_bg(slide, C_DARK)

    # 装饰色块
    add_rect(slide, Inches(0), Inches(4.8), W, Inches(2.7), fill=C_MID)
    add_rect(slide, Inches(0), Inches(6.5), W, Inches(1.0), fill=C_ACCENT)

    # 主标题
    add_text(slide, "基于人声哼唱的旋律转译与风格迁移系统",
             Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.4),
             size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 副标题
    add_text(slide, "中  期  答  辩",
             Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.8),
             size=Pt(26), color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    # 姓名 + 日期
    add_text(slide, "徐亦轲　　2026 年 4 月 20 日",
             Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.7),
             size=Pt(20), color=C_WHITE, align=PP_ALIGN.CENTER)

    # 模块标签
    add_text(slide, "负责模块：音频处理 · 容错量化",
             Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6),
             size=Pt(16), color=C_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 2 — 系统架构
# ════════════════════════════════════════════════════════
def slide_arch(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "系统架构", "System Overview")

    modules = [
        ("哼唱输入\nWAV / MP3", C_GRAY,  False),
        ("音频处理\nCREPE F0",  C_MID,   True),
        ("容错量化\nBiLSTM-CRF", C_ACCENT, True),
        ("风格迁移\nVQ-VAE",    C_MID,   False),
        ("音频渲染\nFluidSynth", C_GRAY,  False),
    ]

    box_w = Inches(2.1)
    box_h = Inches(1.4)
    gap   = Inches(0.18)
    total = len(modules) * float(box_w) + (len(modules) - 1) * float(gap)
    start_x = (float(W) - total) / 2
    y = Inches(2.2)

    for i, (label, color, highlight) in enumerate(modules):
        x = start_x + i * (float(box_w) + float(gap))
        # 高亮模块加阴影效果
        if highlight:
            add_rect(slide, x + Inches(0.04), y + Inches(0.04),
                     box_w, box_h, fill=RGBColor(0xCC, 0x60, 0x10))
        rect = add_rect(slide, x, y, box_w, box_h, fill=color)
        add_text(slide, label,
                 x, y + Inches(0.2), box_w, box_h,
                 size=Pt(15), bold=highlight, color=C_WHITE,
                 align=PP_ALIGN.CENTER)

        # 箭头（除最后一个）
        if i < len(modules) - 1:
            ax = x + float(box_w) + Inches(0.02)
            ay = y + float(box_h) / 2 - Inches(0.12)
            add_text(slide, "▶", ax, ay, gap, Inches(0.3),
                     size=Pt(14), color=C_MID, align=PP_ALIGN.CENTER)

    # 数据流标注
    io_labels = [
        (start_x, "WAV 录音"),
        (start_x + float(box_w) + float(gap), "音高序列\n(time, freq, conf)"),
        (start_x + 2*(float(box_w)+float(gap)), "PrettyMIDI\n单轨旋律"),
        (start_x + 3*(float(box_w)+float(gap)), "PrettyMIDI\n旋律+伴奏"),
        (start_x + 4*(float(box_w)+float(gap)), "WAV 音频"),
    ]
    for x, label in io_labels:
        add_text(slide, label, x, y + box_h + Inches(0.15),
                 box_w, Inches(0.6),
                 size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)

    # 高亮说明
    add_rect(slide, Inches(0.3), Inches(5.4), Inches(5.5), Inches(0.55),
             fill=RGBColor(0xFD, 0xF0, 0xE6), line=C_ACCENT, line_w=Pt(1.5))
    add_text(slide, "  ■  橙色为本报告重点模块（音频处理 + 容错量化）",
             Inches(0.4), Inches(5.42), Inches(5.3), Inches(0.5),
             size=Pt(13), color=C_ACCENT)

    add_text(slide, "支持风格：Pop · Jazz · Classical · Folk",
             Inches(6.5), Inches(5.42), Inches(6.5), Inches(0.5),
             size=Pt(13), color=C_GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════
# Slide 3 — 问题定义
# ════════════════════════════════════════════════════════
def slide_problem(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "核心挑战：哼唱音高 → 精确音符", "Problem Definition")

    challenges = [
        ("🎵", "音高不稳定", "哼唱存在颤音、滑音、音高抖动，难以直接映射到离散音符"),
        ("🎼", "八度偏差",   "演唱者习惯高/低八度哼唱，pyin 提取结果与 GT 相差 12 半音"),
        ("⏱", "时间对齐",   "GT 标注为乐谱理想时间，与实际发声存在 200–400ms 系统性偏移"),
        ("⚖", "类别不平衡", "BIO 标注中音符起始帧（B）仅占 ~1.2%，直接训练模型退化为全 I 预测"),
    ]

    for i, (icon, title, desc) in enumerate(challenges):
        x = Inches(0.5) if i % 2 == 0 else Inches(6.8)
        y = Inches(1.5) if i < 2 else Inches(4.0)
        # 卡片背景
        add_rect(slide, x, y, Inches(6.0), Inches(1.9),
                 fill=C_WHITE, line=C_MID, line_w=Pt(1.2))
        # 图标 + 标题
        add_text(slide, f"{icon}  {title}",
                 x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.55),
                 size=Pt(18), bold=True, color=C_DARK)
        # 描述
        add_text(slide, desc,
                 x + Inches(0.25), y + Inches(0.65), Inches(5.5), Inches(1.1),
                 size=Pt(14), color=C_GRAY)


# ════════════════════════════════════════════════════════
# Slide 4 — 方案设计
# ════════════════════════════════════════════════════════
def slide_design(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "方案设计：BiLSTM-CRF 序列标注", "Solution Design")

    # 左侧：模型结构图
    lx = Inches(0.4)
    layers = [
        ("输入特征 (4维)", C_LIGHT,  C_DARK),
        ("BiLSTM × 2层\nhidden=128", C_MID, C_WHITE),
        ("全连接层 → 发射矩阵", C_MID, C_WHITE),
        ("CRF 解码层", C_ACCENT, C_WHITE),
        ("BIO 标签序列", C_GREEN, C_WHITE),
    ]
    bw, bh, bgap = Inches(4.6), Inches(0.72), Inches(0.18)
    ly_start = Inches(1.4)
    for i, (label, bg, fg) in enumerate(layers):
        by = ly_start + i * (bh + bgap)
        add_rect(slide, lx, by, bw, bh, fill=bg, line=C_MID, line_w=Pt(0.8))
        add_text(slide, label, lx, by + Inches(0.1), bw, bh,
                 size=Pt(14), bold=True, color=fg, align=PP_ALIGN.CENTER)
        if i < len(layers) - 1:
            add_text(slide, "↓",
                     lx + bw/2 - Inches(0.2),
                     by + bh,
                     Inches(0.4), bgap + Inches(0.05),
                     size=Pt(14), color=C_MID, align=PP_ALIGN.CENTER)

    # 右侧：关键技巧
    rx = Inches(5.5)
    add_text(slide, "关键工程技巧",
             rx, Inches(1.35), Inches(7.5), Inches(0.55),
             size=Pt(20), bold=True, color=C_DARK)

    tricks = [
        ("① 输入特征",
         "CREPE 音高 + voiced 概率 + 置信度 + BPM\n→ 替换 pyin，音高准确率从 ~62% 提升至 ~80%"),
        ("② 类别不平衡",
         "B 类仅占 1.2%，加权 CE 辅助损失（B×50）\n+ CRF loss 权重降至 0.2，防止模型崩溃"),
        ("③ 八度修正",
         "比较 pyin 中位音高与 GT 中位音高\n自动吸附至最近 12 半音整数倍补偿"),
        ("④ 时间对齐",
         "互相关算法估算 GT 与实际发声全局偏移\n限制在 ±1s，per-sample 修正标签"),
    ]

    for i, (t, d) in enumerate(tricks):
        ty = Inches(2.0) + i * Inches(1.3)
        add_rect(slide, rx, ty, Inches(7.5), Inches(1.15),
                 fill=C_WHITE, line=C_LIGHT, line_w=Pt(1.0))
        add_rect(slide, rx, ty, Inches(0.08), Inches(1.15), fill=C_ACCENT)
        add_text(slide, t,
                 rx + Inches(0.18), ty + Inches(0.05), Inches(7.2), Inches(0.38),
                 size=Pt(14), bold=True, color=C_DARK)
        add_text(slide, d,
                 rx + Inches(0.18), ty + Inches(0.42), Inches(7.2), Inches(0.7),
                 size=Pt(12), color=C_GRAY)


# ════════════════════════════════════════════════════════
# Slide 5 — 实验结果
# ════════════════════════════════════════════════════════
def make_results_chart() -> bytes:
    methods = ["Baseline\n(pyin)", "v1\n(pyin,2k)", "v4\n(CREPE,13k)", "v5\n(CREPE+对齐)"]
    note_acc = [0.308, 0.237, 0.342, 0.510]
    precision = [0.154, 0.233, 0.359, 0.417]
    f1        = [0.199, 0.234, 0.350, 0.453]

    x = np.arange(len(methods))
    w = 0.26

    fig, ax = plt.subplots(figsize=(9, 4.5))
    fig.patch.set_facecolor("#F5F8FC")
    ax.set_facecolor("#F5F8FC")

    bars1 = ax.bar(x - w, note_acc, w, label="Note Accuracy", color="#2E75B6", alpha=0.9)
    bars2 = ax.bar(x,     precision, w, label="Precision",     color="#70AD47", alpha=0.9)
    bars3 = ax.bar(x + w, f1,        w, label="F1",            color="#ED7D31", alpha=0.9)

    # 最后一组（v5）加高亮边框
    for bar in [bars1[3], bars2[3], bars3[3]]:
        bar.set_edgecolor("#C00000")
        bar.set_linewidth(2.0)

    # 数值标注
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
                    f"{bar.get_height():.3f}",
                    ha="center", va="bottom", fontsize=8.5, fontweight="bold")

    # 目标线
    ax.axhline(0.308 + 0.10, color="gray", linestyle="--", linewidth=1.2, alpha=0.7)
    ax.text(3.45, 0.308 + 0.10 + 0.008, "目标线 (baseline+10%)",
            ha="right", fontsize=8, color="gray")

    ax.set_xticks(x)
    ax.set_xticklabels(methods, fontsize=10)
    ax.set_ylim(0, 0.65)
    ax.set_ylabel("Score", fontsize=11)
    ax.legend(fontsize=10, loc="upper left")
    ax.grid(axis="y", alpha=0.3, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close()
    return buf.getvalue()


def slide_results(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "实验结果（HumTrans TEST 集，769 条）", "Experimental Results")

    # 插入图表
    chart_data = make_results_chart()
    chart_io = io.BytesIO(chart_data)
    slide.shapes.add_picture(chart_io, Inches(0.3), Inches(1.15), Inches(8.2), Inches(4.6))

    # 右侧关键结论
    rx = Inches(8.7)
    add_text(slide, "关键结论",
             rx, Inches(1.3), Inches(4.3), Inches(0.5),
             size=Pt(18), bold=True, color=C_DARK)

    conclusions = [
        (C_GREEN,  "v5 Note Accuracy\n0.510  (+20.1% vs Baseline)"),
        (C_GREEN,  "F1 Score\n0.453  (+127% vs Baseline)"),
        (C_MID,    "CREPE 替换 pyin\n+3.5% Note Accuracy"),
        (C_MID,    "时间对齐 + 八度修正\n大幅改善标注质量"),
        (C_ACCENT, "目标 ≥ +10% ✓ 达标"),
    ]

    for i, (color, text) in enumerate(conclusions):
        cy = Inches(1.95) + i * Inches(1.02)
        add_rect(slide, rx, cy, Inches(4.3), Inches(0.85),
                 fill=C_WHITE, line=color, line_w=Pt(2.0))
        add_rect(slide, rx, cy, Inches(0.09), Inches(0.85), fill=color)
        add_text(slide, text,
                 rx + Inches(0.18), cy + Inches(0.05), Inches(4.0), Inches(0.78),
                 size=Pt(13), bold=(i == 4), color=C_DARK if i < 4 else color)


# ════════════════════════════════════════════════════════
# Slide 6 — 可视化 Demo
# ════════════════════════════════════════════════════════
def slide_demo(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "可视化演示工具", "Demo")

    # 左侧：功能说明
    features = [
        ("🎹  钢琴卷帘对比",    "GT 标注 vs BiLSTM-CRF 预测，直观展示音符准确率"),
        ("📊  B 类得分曲线",     "模型对每帧「音符起始」的置信度，直观展示检测逻辑"),
        ("🔊  音频对比播放",     "原始哼唱 / Baseline 量化 / BiLSTM-CRF 量化，三路对比"),
        ("⚙️  参数交互调节",     "实时调整 peak 阈值和 distance，观察对音符分割的影响"),
    ]

    lx = Inches(0.4)
    for i, (title, desc) in enumerate(features):
        fy = Inches(1.4) + i * Inches(1.35)
        add_rect(slide, lx, fy, Inches(5.8), Inches(1.15),
                 fill=C_WHITE, line=C_MID, line_w=Pt(1.0))
        add_rect(slide, lx, fy, Inches(0.08), Inches(1.15), fill=C_MID)
        add_text(slide, title,
                 lx + Inches(0.18), fy + Inches(0.05), Inches(5.5), Inches(0.4),
                 size=Pt(15), bold=True, color=C_DARK)
        add_text(slide, desc,
                 lx + Inches(0.18), fy + Inches(0.48), Inches(5.5), Inches(0.6),
                 size=Pt(12), color=C_GRAY)

    # 右侧：启动说明框
    rx = Inches(6.5)
    add_rect(slide, rx, Inches(1.4), Inches(6.6), Inches(2.5),
             fill=RGBColor(0x1F, 0x38, 0x64), line=None)
    add_text(slide, "启动命令",
             rx + Inches(0.2), Inches(1.5), Inches(6.2), Inches(0.5),
             size=Pt(15), bold=True, color=C_LIGHT)
    add_text(slide,
             "streamlit run tools/visualizer.py\n  --server.port 8501",
             rx + Inches(0.2), Inches(2.05), Inches(6.2), Inches(0.75),
             size=Pt(13), color=RGBColor(0xA8, 0xD8, 0xFF))
    add_text(slide, "浏览器访问 http://localhost:8501",
             rx + Inches(0.2), Inches(2.85), Inches(6.2), Inches(0.45),
             size=Pt(12), color=C_LIGHT, italic=True)

    # 截图占位框
    add_rect(slide, rx, Inches(4.05), Inches(6.6), Inches(2.85),
             fill=C_WHITE, line=C_MID, line_w=Pt(1.5))
    add_text(slide, "【此处粘贴 Streamlit 截图】",
             rx, Inches(4.05), Inches(6.6), Inches(2.85),
             size=Pt(16), color=C_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 7 — 总结
# ════════════════════════════════════════════════════════
def slide_summary(prs):
    slide = blank_slide(prs)
    slide_bg(slide)
    header_bar(slide, "小结与展望", "Summary")

    # 左：已完成
    add_rect(slide, Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.5),
             fill=C_MID)
    add_text(slide, "  ✅  已完成",
             Inches(0.4), Inches(1.3), Inches(5.9), Inches(0.5),
             size=Pt(17), bold=True, color=C_WHITE)

    done = [
        "CREPE F0 特征提取管线（含 pyin 自动回退）",
        "BiLSTM-CRF 序列标注量化器（533K 参数）",
        "八度偏移自动修正 + 互相关时间对齐",
        "HumTrans 全量训练（13k 样本，50 epoch）",
        "Note Accuracy 0.510，超越 Baseline +20.1%  ✓",
        "Streamlit 交互可视化工具",
    ]
    for i, item in enumerate(done):
        dy = Inches(1.95) + i * Inches(0.67)
        color = C_GREEN if i == 4 else C_DARK
        bold  = (i == 4)
        add_text(slide, f"  •  {item}",
                 Inches(0.5), dy, Inches(5.7), Inches(0.62),
                 size=Pt(14), bold=bold, color=color)

    # 右：下一步
    add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.9), Inches(0.5),
             fill=C_ACCENT)
    add_text(slide, "  🔜  下一步计划",
             Inches(7.0), Inches(1.3), Inches(5.9), Inches(0.5),
             size=Pt(17), bold=True, color=C_WHITE)

    next_steps = [
        "接入风格迁移模块，完成端到端 pipeline 联调",
        "录制完整 Demo 视频（哼唱 → 多风格输出）",
        "优化推理速度，达到实时响应（< 2s）",
        "（可选）探索 Transformer 替换 LSTM",
    ]
    for i, item in enumerate(next_steps):
        ny = Inches(1.95) + i * Inches(0.67)
        add_text(slide, f"  •  {item}",
                 Inches(7.1), ny, Inches(5.7), Inches(0.62),
                 size=Pt(14), color=C_DARK)

    # 底部一句话
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.72),
             fill=C_DARK)
    add_text(slide,
             "BiLSTM-CRF 量化器已达成中期目标，为后续端到端系统打下坚实基础",
             Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.72),
             size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# 主程序
# ════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_arch(prs)
    slide_problem(prs)
    slide_design(prs)
    slide_results(prs)
    slide_demo(prs)
    slide_summary(prs)

    out = Path("reports/midterm_xyk.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"已生成：{out.resolve()}")


if __name__ == "__main__":
    main()
